from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def add_content_slide(presentation: Presentation, title: str, content: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content


def add_chart_slide(presentation: Presentation, title: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = title

    chart_data = CategoryChartData()
    chart_data.categories = ["08:45", "09:00", "09:15", "09:30"]
    chart_data.add_series("Sample value", (10000, 10120, 9980, 10050))

    slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(2),
        Inches(2),
        Inches(5),
        Inches(3),
        chart_data,
    )


def create_sample_presentation(output_pptx: Path) -> None:
    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Math Teaching Material Demo"
    title_slide.placeholders[1].text = "Generated with python-pptx"

    add_content_slide(
        presentation,
        "Teaching Idea",
        "Use code to generate repeatable classroom materials.",
    )
    add_content_slide(
        presentation,
        "Workflow",
        "Prepare data, generate slides, then review and polish by hand.",
    )
    add_chart_slide(presentation, "Sample Line Chart")

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_pptx)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate a sample PowerPoint deck.")
    parser.add_argument("output_pptx", type=Path)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    create_sample_presentation(args.output_pptx)
    print(f"Wrote PowerPoint file: {args.output_pptx}")


if __name__ == "__main__":
    main()

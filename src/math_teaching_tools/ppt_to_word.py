from __future__ import annotations

import argparse
from pathlib import Path

from docx import Document
from pptx import Presentation


def ppt_to_word(ppt_file: Path, word_file: Path) -> None:
    presentation = Presentation(ppt_file)
    document = Document()

    for slide_number, slide in enumerate(presentation.slides, start=1):
        document.add_heading(f"Slide {slide_number}", level=1)

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            text = shape.text.strip()
            if text:
                document.add_paragraph(text)

        if slide_number < len(presentation.slides):
            document.add_page_break()

    word_file.parent.mkdir(parents=True, exist_ok=True)
    document.save(word_file)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Extract text from a PowerPoint deck into a Word file.",
    )
    parser.add_argument("ppt_file", type=Path)
    parser.add_argument("word_file", type=Path)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    ppt_to_word(args.ppt_file, args.word_file)
    print(f"Wrote Word file: {args.word_file}")


if __name__ == "__main__":
    main()
